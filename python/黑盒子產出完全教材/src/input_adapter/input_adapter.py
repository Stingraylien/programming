import os  # 匯入內建模組以處理系統路徑操作
from typing import Dict, List  # 匯入型態標註輔助功能以利型態檢查

def process_file(str_file_path: str) -> Dict:  # 定義接收實體檔案路徑並回傳標準化中間格式字典的主入口函式
    """Process a single file and return intermediate data structure."""  # 函式說明文件字串，描述此功能為處理單一檔案並轉化為中間結構
    if not os.path.exists(str_file_path):  # 邏輯判斷，檢查傳入的實體檔案路徑是否真的存在於作業系統中
        raise FileNotFoundError(f"File not found: {str_file_path}")  # 例外拋出，若檔案不存在則中斷執行並顯示檔案未找到之錯誤
    
    str_ext = os.path.splitext(str_file_path)[1].lower()  # 字串變數，拆解檔名以取得最後段的副檔名，並轉為小寫以利後續邏輯比對

    if str_ext == '.pdf':  # 條件分支，判斷檔案類型是否為可攜式文件格式 (PDF)
        return process_pdf(str_file_path)  # 回傳值，呼叫並執行專供 PDF 套件解析的子函式結果
    elif str_ext == '.pptx':  # 條件分支，判斷檔案是否為 PowerPoint 2007+ 版簡報格式
        return process_pptx(str_file_path)  # 回傳值，呼叫專屬處理簡報投影片內容的轉譯函式
    elif str_ext == '.ppt':  # 條件分支，判斷檔案是否為 PowerPoint 97-2003 舊版簡報格式
        return process_ppt(str_file_path)  # 回傳值，嘗試透過 LibreOffice 轉換後解析，或回傳指引
    elif str_ext == '.docx':  # 條件分支，判斷檔案是否為 Word 2007+ 版文件格式
        return process_docx(str_file_path)  # 回傳值，呼叫專責處理 Word 段落內容的轉譯函式
    elif str_ext == '.doc':  # 條件分支，判斷檔案是否為 Word 97-2003 舊版文件格式
        return process_doc(str_file_path)  # 回傳值，嘗試透過 LibreOffice 轉換後解析，或回傳指引
    elif str_ext == '.mp4':  # 條件分支，判斷檔案是否為常見的 H.264 數位影片格式
        return process_mp4(str_file_path)  # 回傳值，呼叫目前預留作為影片音訊转成文字的預處理函式
    elif str_ext == '.txt':  # 條件分支，判斷檔案是否為最基礎的 UTF-8 純文字格式
        return process_txt(str_file_path)  # 回傳值，呼叫處理純文字切割與內文讀取的函式
    else:  # 當檔案副檔名不在上述預先定義的受支援白名單清單中時
        raise ValueError(f"Unsupported file type: {str_ext}")  # 例外拋出，告知使用者此檔案格式目前尚未受系統支援開發

def process_pdf(str_file_path: str) -> Dict:  # 定義專門用來將 PDF 內容轉譯為系統中間字典格式的子函式
    """Process PDF file."""  # 函式說明文件字串，標記此區塊為 PDF 檔案解析主邏輯
    try:  # 啟動保護機制，避免因第三方套件缺失或 PDF 格式異常導致整個 API 服務崩潰
        import PyPDF2  # 模組引入，在程式運行時嘗試動態載入 PDF 處理核心套件 PyPDF2
        list_sections = []  # 陣列型字典變數，初始化一個空白清單用來存放從每一頁提取出來的標題與內容
        with open(str_file_path, 'rb') as obj_file:  # 物件變數，以二進位唯讀模式開啟 PDF 實體檔案並建立物件實例
            obj_pdf_reader = PyPDF2.PdfReader(obj_file)  # 物件變數，利用套件功能建立一個可供程式逐頁檢視的閱讀器物件
            for int_page_num, obj_page in enumerate(obj_pdf_reader.pages):  # 整數與物件迴圈，從第一頁開始編號並走訪整本 PDF 內容
                list_sections.append({  # 在區段清單中加入一個代表該頁面的封裝字典物件
                    "title": f"Page {int_page_num + 1}",  # 字串變數，將頁碼加一後格式化為讀者能理解的頁次標題
                    "content": obj_page.extract_text()  # 字串變數，透過掃描功能從目前頁面中抽出原始的純文字內容資料
                })  # 結束該頁面字典的添加動作
        dict_result = {  # 字典變數，建立最終要拋回給核心引擎的標準中間格式解析成果字典
            "source_type": "pdf",  # 字串標記，註明這份內容最初是來自 PDF 格式的轉換
            "sections": list_sections if list_sections else [{"title": "PDF Content", "content": "Sample"}]  # 陣列清單，若有內容則轉移，否則給予預設占位資料
        }  # 結束成果字典的宣告
        return dict_result  # 回傳值，傳回這份完整封裝好的解析成果給主控流程
    except ImportError:  # 錯誤捕捉，當伺服器環境完全沒有安裝 PyPDF2 所需之相依套件時
        dict_fallback = {  # 字典變數，建立一份包含錯誤指示說明的退化型成果包
            "source_type": "pdf",  # 字串標記，指出此解析原本的目標是 PDF 格式
            "sections": [{  # 陣列清單，置入一個帶有錯誤警示說明的虛擬章節內容
                "title": "PDF Document Error",  # 字串變數，顯示出錯的標題指示
                "content": "PDF requires PyPDF2. Install with: pip install PyPDF2"  # 字串變數，提供白話解方要求使用者安裝套件
            }]  # 結束虛擬章節
        }  # 結束退化字典定義
        return dict_fallback  # 回傳值，將包含錯誤提示的字典回傳以供前端顯示

def process_pptx(str_file_path: str) -> Dict:  # 定義專責處理 PowerPoint 投影片解析任務的函式
    """Process PPTX file."""  # 函式說明文件字串，標識此功能的處理目標為簡報檔
    try:  # 啟動保護，防止載入簡報或抽取文字時發生非預期之中斷
        from pptx import Presentation  # 模組引入，嘗試動態從系統庫中叫用 PowerPoint 解析的核心元件
        list_sections = []  # 陣列型字典變數，初始化緩衝清單準備接收每一張投影片的精華摘要
        obj_prs = Presentation(str_file_path)  # 物件變數，將實體簡報檔案映射到記憶體中的物件模型以供遍歷
        for int_slide_num, obj_slide in enumerate(obj_prs.slides):  # 整數與物件迴圈，逐一依序存取簡報中的所有投影片內容
            list_slide_content = []  # 陣列字串變數，建立臨時清單來存放該頁面上所有零碎的文字框碎片
            for obj_shape in obj_slide.shapes:  # 物件迴圈，掃描目前投影片上所有的幾何圖形、文字框與多媒體組件
                if hasattr(obj_shape, "text"):  # 邏輯判斷，只有當該組件確定具備可讀取的文字屬性時才進行處理
                    list_slide_content.append(obj_shape.text)  # 字串追加，將讀取到的文字方塊內容存進當頁的緩衝清單中
            
            list_sections.append({  # 將該張投影片彙整出的所有資訊打包成一個單獨的小節字典
                "title": f"Slide {int_slide_num + 1}",  # 字串變數，將投影片索引編號美化為人類可讀的頁次標題
                "content": "\n".join(list_slide_content)  # 字串變數，將同頁的多個文字碎片用換行符號黏合成一段完整文章
            })  # 結束投影片字典打包
        
        dict_result = {  # 字典變數，定義最後輸出給系統讀取的轉譯成果
            "source_type": "pptx",  # 字串標記，定義來源型態為簡報檔格式
            "sections": list_sections if list_sections else [{"title": "PPTX", "content": "Empty"}]  # 陣列清單，承接所有投影片內容
        }  # 結束字典定義
        return dict_result  # 回傳值，交出解析好的簡報內容字典
    except ImportError:  # 錯誤捕捉，當環境缺乏簡報解析必要的 python-pptx 庫時
        dict_fallback = {  # 字典型變數，打包錯誤訊息回傳
            "source_type": "pptx",  # 標記來源
            "sections": [{  # 建置錯誤指示章節
                "title": "PowerPoint Error",  # 出錯標題
                "content": "PPTX requires python-pptx. Install with: pip install python-pptx"  # 白話說明修復方法
            }]  # 結束
        }  # 結束
        return dict_fallback  # 退回錯誤字典

def process_docx(str_file_path: str) -> Dict:  # 定義專責處理 Word 文件段落與標題解析的專業函式
    """Process DOCX file."""  # 函式說明文件字串，指定此功能專修 DOCX 格式
    try:  # 啟動例外捕捉，防範損毀的 Word 檔導致崩毀
        from docx import Document  # 模組引入，引進專門拆解 OpenXML 文件格式的官方分析器
        obj_doc = Document(str_file_path)  # 物件變數，建立一個指向該 Word 檔案的記憶體實例以利檢索
        list_sections = []  # 陣列字典變數，準備蒐集根據標題切割出來的文件分段清單
        
        str_current_section_title = "Document"  # 字串變數，預設一個標題稱號，後續將被真實標題替代
        list_current_content = []  # 陣列字串變數，用來積存屬於同一個標題下方的所有內文段落
        
        for obj_para in obj_doc.paragraphs:  # 物件迴圈，依序走過 Word 檔案中從頭到尾的所有自然段落
            if obj_para.style.name.startswith('Heading'):  # 邏輯判斷，若偵測到當前段落被設定為 Word 系統內建的標題樣式
                if list_current_content:  # 邏輯判斷，若先前已經有積攢了一部分內文素材
                    list_sections.append({  # 將先前的素材連同其標題一起打包存入成果清單中
                        "title": str_current_section_title,  # 字串變數，取用先前的標題作為章節名稱
                        "content": "\n\n".join(list_current_content)  # 字串變數，用雙換行將清單內容糊合成一篇短文
                    })  # 結束章節添加
                str_current_section_title = obj_para.text  # 字串變數，將當前讀到的新標題文字設定為未來的章節主旨
                list_current_content = []  # 清空清單，重新開始收集下方新的內文段落
            else:  # 當讀到的是平鋪直敘的一般內容段落時
                if obj_para.text.strip():  # 邏輯判斷，只要這行不是純粹的隱形空白字元
                    list_current_content.append(obj_para.text)  # 字串追加，將此行內文存入目前所屬的章節內容桶子裡
        
        if list_current_content:  # 邊際處理，當檔案掃描完畢後若還有沒存進清單的尾巴內容
            list_sections.append({  # 強制將最終殘留的一截內容也進行存檔打包
                "title": str_current_section_title,  # 沿用最後一個偵測到的標題
                "content": "\n\n".join(list_current_content)  # 打包黏合
            })  # 結束最後打包
        
        dict_result = {  # 字典變數，宣告最終解析成果
            "source_type": "docx",  # 標註來源為 Word
            "sections": list_sections if list_sections else [{"title": "Word", "content": "Sample"}]  # 承接所有段落
        }  # 結尾
        return dict_result  # 交出轉譯成功的字典
    except ImportError:  # 處理依賴套件缺失之狀況
        dict_fallback = {  # 字典變數，建置錯誤封裝
            "source_type": "docx",  # 標明格式
            "sections": [{  # 單一錯誤節
                "title": "Word Error",  # 標題
                "content": "DOCX requires python-docx. Install with: pip install python-docx"  # 說明
            }]  # 結尾
        }  # 結束
        return dict_fallback  # 回傳錯誤

def _try_libreoffice_convert(str_file_path: str, str_target_ext: str) -> str | None:
    """嘗試用 LibreOffice 將舊版 Office 格式轉為新格式，回傳轉換後路徑或 None"""
    import subprocess
    str_out_dir = os.path.dirname(str_file_path)
    try:
        result = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', str_target_ext, '--outdir', str_out_dir, str_file_path],
            capture_output=True, timeout=60
        )
        if result.returncode == 0:
            str_converted = os.path.splitext(str_file_path)[0] + '.' + str_target_ext
            if os.path.exists(str_converted):
                return str_converted
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass
    return None

def process_ppt(str_file_path: str) -> Dict:
    """Process old-format PPT file by converting via LibreOffice, then delegating to process_pptx."""
    str_converted = _try_libreoffice_convert(str_file_path, 'pptx')
    if str_converted:
        return process_pptx(str_converted)
    return {
        "source_type": "ppt",
        "sections": [{
            "title": "PowerPoint 97-2003 (.ppt)",
            "content": (
                "舊版 .ppt 格式需要 LibreOffice 才能轉換解析。"
                "請安裝 LibreOffice (https://www.libreoffice.org) 或將檔案另存為 .pptx 後重新上傳。"
            )
        }]
    }

def process_doc(str_file_path: str) -> Dict:
    """Process old-format DOC file by converting via LibreOffice, then delegating to process_docx."""
    str_converted = _try_libreoffice_convert(str_file_path, 'docx')
    if str_converted:
        return process_docx(str_converted)
    return {
        "source_type": "doc",
        "sections": [{
            "title": "Word 97-2003 (.doc)",
            "content": (
                "舊版 .doc 格式需要 LibreOffice 才能轉換解析。"
                "請安裝 LibreOffice (https://www.libreoffice.org) 或將檔案另存為 .docx 後重新上傳。"
            )
        }]
    }

def process_mp4(str_file_path: str) -> Dict:  # 定義影片處理預設函式，目前作為未來語音轉文字擴充的樁位
    """Process MP4 file."""  # 函式說明，註明此處處理影片格式
    dict_placeholder = {  # 字典型變數，建立一個預留未來真正實作後要回傳的內容骨架
        "source_type": "mp4",  # 字串變數，標註來源是影片格式
        "sections": [{  # 陣列清單，放入一段目前的實太提示
            "title": "Video Instruction",  # 標題文字
            "content": "Video processing via moviepy is supported in the video pipeline."  # 字串變數，說明目前的整合現況
        }]  # 結尾
    }  # 結束
    return dict_placeholder  # 回傳目前暫用之占位字典

def process_txt(str_file_path: str) -> Dict:  # 定義專責處理最原始、最常見之純文字 UTF-8 格式的解析函式
    """Process TXT file."""  # 函式說明文件字串，指定此處負責 TXT 分析
    try:  # 啟動保護，防止檔案開啟失敗或編碼錯誤
        list_sections = []  # 陣列字典變數，初始化用來存放依照特定符號切分後的小節
        with open(str_file_path, 'r', encoding='utf-8') as obj_file:  # 物件變數，以 UTF-8 編碼開啟並建立文件讀取控制流
            str_content = obj_file.read()  # 字串變數，毫不保留地一次性將整個檔案的文字內容抓進記憶體變數中
        
        if '###' in str_content:  # 邏輯判斷，檢查文字內是否含有系統建議的章節分界標記
            list_parts = str_content.split('###')  # 陣列字串變數，沿著分界標記將全本內容暴力拆解成多個小塊
            for int_i, str_part in enumerate(list_parts):  # 整數與字串迴圈，對切出來的每一塊文字碎片進行加工處理
                str_part = str_part.strip()  # 字串處理，自動修剪掉該碎片前後剩餘的無謂空白與換行殘渣
                if str_part:  # 邏輯判斷，只有該碎片內含實質文字訊息時才錄用
                    list_lines = str_part.split('\n')  # 陣列字串變數，再次將此碎片延換行符號切成更細的單行集合
                    str_title = list_lines[0].strip() if list_lines else f"Section {int_i}"  # 字串變數，嘗試以第一行文字做為小節標題
                    str_body = '\n'.join(list_lines[1:]).strip() if len(list_lines) > 1 else ''  # 字串變數，將剩餘行數糊死做為本段真實內文
                    list_sections.append({  # 將目前萃取出的標題與內文封裝至這組章節字典中
                        "title": str_title,  # 置入標題
                        "content": str_body  # 置入內文
                    })  # 結束添加
        else:  # 若該文件只是一整篇平鋪直敘、沒有特殊切分記號的普通文字檔
            list_sections = [{  # 陣列變數，將檔案視同為單一巨大章節直接進行封裝
                "title": os.path.splitext(os.path.basename(str_file_path))[0],  # 字串變數，將檔案的純檔名擷取出來作為唯一的預定標題
                "content": str_content  # 字串變數，完整承接原本的所有文字作為內容
            }]  # 結束唯一字典定義
        
        dict_result = {  # 字典變數，宣告給前端回看的最終解析結果包
            "source_type": "txt",  # 標記身分為文字檔
            "sections": list_sections if list_sections else [{"title": "Text", "content": "Empty"}]  # 承載解析出的內容或提示
        }  # 結尾
        return dict_result  # 交出轉錄好的內容字典
    except Exception as obj_err:  # 錯誤捕捉機制，對付那些連檔案都打不開的悲慘慘案
        dict_err = {  # 字典變數，封裝含有電腦底層錯誤描述的警示字典
            "source_type": "txt",  # 品種標記
            "sections": [{  # 錯誤指示章節
                "title": "Text File Access Error",  # 錯誤名稱
                "content": f"System error details: {str(obj_err)}"  # 將編譯器噴出的英文錯誤原因包在內文中說明
            }]  # 結束
        }  # 結束字典定義
        return dict_err  # 打退堂鼓，回傳錯誤報告