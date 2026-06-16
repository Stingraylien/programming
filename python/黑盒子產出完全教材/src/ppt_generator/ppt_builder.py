from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors — matches cover_builder.py
_NAVY = RGBColor(18, 52, 99)
_ACCENT = RGBColor(60, 140, 230)
_WHITE = RGBColor(255, 255, 255)
_LIGHT_BG = RGBColor(240, 245, 255)
_BODY_TEXT = RGBColor(30, 50, 90)
_SUB_TEXT = RGBColor(80, 110, 160)

_INT_BLANK_LAYOUT = 6  # blank slide layout index in default template


def _set_bg(slide, rgb_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color


def _add_rect(slide, left, top, width, height, rgb_fill):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_fill
    shape.line.fill.background()
    return shape


def _add_label(slide, str_text, left, top, width, height,
               int_pt, bool_bold=False, rgb_color=_WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str_text
    run.font.size = Pt(int_pt)
    run.font.bold = bool_bold
    run.font.color.rgb = rgb_color
    return txBox


def _build_title_slide(obj_prs, str_course_name):
    obj_slide = obj_prs.slides.add_slide(obj_prs.slide_layouts[_INT_BLANK_LAYOUT])
    _set_bg(obj_slide, _NAVY)
    _add_rect(obj_slide, Inches(0), Inches(0), Inches(13.33), Inches(0.12), _ACCENT)
    _add_rect(obj_slide, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12), _ACCENT)
    _add_label(obj_slide, str_course_name,
               Inches(1), Inches(2.3), Inches(11.33), Inches(1.9),
               int_pt=44, bool_bold=True, rgb_color=_WHITE, align=PP_ALIGN.CENTER)
    _add_label(obj_slide, "課程重點摘要",
               Inches(1), Inches(4.4), Inches(11.33), Inches(0.9),
               int_pt=28, rgb_color=_ACCENT, align=PP_ALIGN.CENTER)


def _build_content_slide(obj_prs, dict_content):
    obj_slide = obj_prs.slides.add_slide(obj_prs.slide_layouts[_INT_BLANK_LAYOUT])
    _set_bg(obj_slide, _LIGHT_BG)

    # Left accent bar
    _add_rect(obj_slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), _NAVY)

    # Chapter title
    str_chapter = dict_content.get('chapter', '章節重點')
    _add_label(obj_slide, str_chapter,
               Inches(0.35), Inches(0.18), Inches(12.8), Inches(0.95),
               int_pt=28, bool_bold=True, rgb_color=_NAVY)

    # Separator line under title
    _add_rect(obj_slide, Inches(0.35), Inches(1.2), Inches(12.63), Inches(0.04), _ACCENT)

    # Body text box
    list_objectives = dict_content.get('learning_objectives', [])
    list_paragraphs = dict_content.get('paragraphs', [])

    txBox = obj_slide.shapes.add_textbox(Inches(0.35), Inches(1.35), Inches(12.63), Inches(5.95))
    tf = txBox.text_frame
    tf.word_wrap = True

    bool_first = True

    for str_obj in list_objectives:
        p = tf.paragraphs[0] if bool_first else tf.add_paragraph()
        bool_first = False
        run = p.add_run()
        run.text = f"• {str_obj}"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = _NAVY

    for str_para in list_paragraphs[:3]:
        p = tf.paragraphs[0] if bool_first else tf.add_paragraph()
        bool_first = False
        str_truncated = str_para[:150] + ("…" if len(str_para) > 150 else "")
        run = p.add_run()
        run.text = f"    {str_truncated}"
        run.font.size = Pt(15)
        run.font.color.rgb = _SUB_TEXT


def build_summary_ppt(str_course_name, list_canonical_contents, str_output_path):
    obj_prs = Presentation()
    obj_prs.slide_width = Inches(13.33)
    obj_prs.slide_height = Inches(7.5)

    _build_title_slide(obj_prs, str_course_name)

    for dict_content in list_canonical_contents:
        _build_content_slide(obj_prs, dict_content)

    obj_prs.save(str_output_path)
    return str_output_path
