# input_adapter.py
# Input Adapter: Convert various material formats to intermediate data

import os  # 匯入內建模組以處理系統路徑操作
from typing import Dict, List  # 匯入型態標註輔助功能以利型態檢查

def process_file(str_file_path: str) -> Dict:  # 定義接收實體路徑字串並回傳字典的主函式
    """
    Process a single file and return intermediate data structure.
    """
    if not os.path.exists(str_file_path):  # 若傳入的該路徑根本不存在系統中
        raise FileNotFoundError(f"File not found: {str_file_path}")  # 直接拋出找不到檔案的致命錯誤例外
    
    str_ext = os.path.splitext(str_file_path)[1].lower()  # 字串變數，切割檔名以取得最尾段的副檔名並轉小寫比對

    if str_ext == '.pdf':  # 判斷是否為 PDF 格式
        return process_pdf(str_file_path)  # 交給專責 PDF 解析函式並回傳結果字典
    elif str_ext == '.pptx':  # 判斷是否為 PPTX 簡報格式
        return process_pptx(str_file_path)  # 委由簡報解析函式處理
    elif str_ext == '.docx':  # 判斷是否為 DOCX 文件檔
        return process_docx(str_file_path)  # 呼叫專屬函式讀取段落
    elif str_ext == '.mp4':  # 判斷是否為數位影片檔
        return process_mp4(str_file_path)  # 派遣影片拆解函式處理
    elif str_ext == '.txt':  # 判斷是否為純文字筆記
        return process_txt(str_file_path)  # 使用預設的文字拆解工具
    else:  # 若都不是我們定義過的白名單副檔名
        raise ValueError(f"Unsupported file type: {str_ext}")  # 拋出數值錯誤，拒絕未知的檔案格式

def process_pdf(str_file_path: str) -> Dict:  # 定義處理 PDF 轉文字的函式
    """Process PDF file."""
    try:  # 建立基礎除錯保護層
        import PyPDF2  # 嘗試動態匯入第三方套件 PyPDF2
        list_sections = []  # 陣列變數，準備存放每頁 PDF 解析結果
        with open(str_file_path, 'rb') as file_obj:  # 使用二進位唯讀模式開啟 PDF
            obj_pdf_reader = PyPDF2.PdfReader(file_obj)  # 物件變數，實例化 PDF 閱讀器
            for int_page_num, obj_page in enumerate(obj_pdf_reader.pages):  # 走訪 PDF 中的所有頁面物件
                list_sections.append({  # 在清單中塞入每頁內容的字典
                    "title": f"Page {int_page_num + 1}",  # 給予預設標題 "Page N"
                    "content": obj_page.extract_text()  # 抽出該頁面的純文字字串
                })
        dict_result = {  # 字典變數，最終輸出的標準化封裝格式
            "source_type": "pdf",  # 標註這份資料來源是 PDF
            "sections": list_sections if list_sections else [{"title": "PDF Content", "content": "Sample PDF content"}]  # 若有抽出內容則回傳，否則塞假資料墊檔
        }
        return dict_result  # 將整理完的字典回傳
    except ImportError:  # 若系統遺失了 PyPDF2 套件
        dict_fallback = {  # 字典變數，退回告知必須安裝套件的提示內容
            "source_type": "pdf",  # 一樣標明他是 PDF
            "sections": [{  # 塞入人工警示訊息
                "title": "PDF Document",  # 標題填入預設名稱
                "content": "PDF processing requires PyPDF2. Please install it with: pip install PyPDF2"  # 告知缺安裝的套件
            }]
        }
        return dict_fallback  # 回傳這包失敗的結果

def process_pptx(str_file_path: str) -> Dict:  # 定義處理 PowerPoint 檔案的核心函式
    """Process PPTX file."""
    try:  # 嘗試載入模組與解析
        from pptx import Presentation  # 動態引入 python-pptx 模組
        list_sections = []  # 陣列變數，用來承裝每一張投影片的文字
        obj_prs = Presentation(str_file_path)  # 物件變數，載入簡報檔
        for int_slide_num, obj_slide in enumerate(obj_prs.slides):  # 走訪每一張投影片
            list_slide_content = []  # 陣列變數，暫存各文字方塊內容
            for obj_shape in obj_slide.shapes:  # 輪詢投影片上的圖形
                if hasattr(obj_shape, "text"):  # 判斷是否具備文字屬性
                    list_slide_content.append(obj_shape.text)  # 取出文字並收集
            
            list_sections.append({  # 打包該投影片所有文字
                "title": f"Slide {int_slide_num + 1}",  # 給予預設 Slide 標題
                "content": "\n".join(list_slide_content)  # 將文字碎片用換行黏合
            })
        
        dict_result = {  # 字典變數，整理回報 JSON 格式
            "source_type": "pptx",  # 標示為簡報
            "sections": list_sections if list_sections else [{"title": "PPTX Content", "content": "Sample PPTX content"}]  # 正式載入內容清單
        }
        return dict_result  # 回傳產出字典
    except ImportError:  # 取不到依賴套件時的防呆
        dict_fallback = {  # 字典變數，回報失敗
            "source_type": "pptx",  # 標題依舊為簡報
            "sections": [{  # 置入提示訊息區塊
                "title": "PowerPoint Presentation",  # 代稱標題
                "content": "PPTX processing requires python-pptx. Please install it with: pip install python-pptx"  # 提示需要安裝的缺失套件
            }]
        }
        return dict_fallback  # 返回例外處置

def process_docx(str_file_path: str) -> Dict:  # 處理 Word 文件的專屬函式
    """Process DOCX file."""
    try:  # 保護主區塊避免當機
        from docx import Document  # 動態抓取套件
        obj_doc = Document(str_file_path)  # 物件變數，讀取記憶體中的 Word 結構
        list_sections = []  # 陣列變數，預計盛裝章節內容
        
        str_current_section_title = "Document"  # 字串變數，用來記憶段落所屬之主標題
        list_current_content = []  # 陣列變數，累積屬於同標題的內文
        
        for obj_para in obj_doc.paragraphs:  # 走訪文件內的所有段落
            if obj_para.style.name.startswith('Heading'):  # 倘若讀到標題樣式
                if list_current_content:  # 若先前已有累積的舊內文
                    list_sections.append({  # 把舊的封裝進清單裡
                        "title": str_current_section_title,  # 取用舊的標題名字
                        "content": "\n\n".join(list_current_content)  # 將陣列黏合成一段文章
                    })
                str_current_section_title = obj_para.text  # 將新取到的標題取代舊的
                list_current_content = []  # 清空陣列準備重新收集下方的內文
            else:  # 若是一般內文段落
                if obj_para.text.strip():  # 只要不是純空白列
                    list_current_content.append(obj_para.text)  # 加入該段落字串
        
        if list_current_content:  # 迴圈結束後，若袋中還殘留未封裝的文字
            list_sections.append({  # 手動封裝最後的內容
                "title": str_current_section_title,  # 繼承最後一任標題
                "content": "\n\n".join(list_current_content)  # 打包黏合
            })
        
        dict_result = {  # 字典變數，標準產出格
            "source_type": "docx",  # 宣告血統為 DOCX
            "sections": list_sections if list_sections else [{"title": "Document Content", "content": "Sample document content"}]  # 順帶給個防呆內容
        }
        return dict_result  # 拋回字典
    except ImportError:  # 處理沒裝套件的狀況
        dict_fallback = {  # 字典變數，失敗包裹
            "source_type": "docx",  # 通知原意需解析 docx
            "sections": [{  # 建構提示
                "title": "Word Document",  # 標題設定
                "content": "DOCX processing requires python-docx. Please install it with: pip install python-docx"  # 要求系統管理員補齊套件
            }]
        }
        return dict_fallback  # 收場回撤

def process_mp4(str_file_path: str) -> Dict:  # 面對影音檔的預留前處理函式
    """Process MP4 file."""
    dict_placeholder = {  # 字典變數，由於實作複雜故暫建空殼
        "source_type": "mp4",  # 標示影片屬性
        "sections": [{  # 陣列內含物件
            "title": "Video Transcript",  # 預定放置字幕之地
            "content": "Video processing requires additional libraries like moviepy or speech_recognition. Please install them with: pip install moviepy speech_recognition"  # 預留要求安裝語音辨識之通知
        }]
    }
    return dict_placeholder  # 歸還預留用字典

def process_txt(str_file_path: str) -> Dict:  # 專門為一般記事本拆解設計的函式
    """Process TXT file."""
    try:  # 防止隨意崩潰
        list_sections = []  # 陣列變數，開空箱準備收資料
        with open(str_file_path, 'r', encoding='utf-8') as file_obj:  # 使用 UTF-8 開啟字串以防亂碼
            str_content = file_obj.read()  # 字串變數，直接通吃全本內容讀入
        
        if '###' in str_content:  # 檢查是否具備俗成的章節切分號
            list_parts = str_content.split('###')  # 陣列變數，根據三個井號暴力切開字串
            for int_i, str_part in enumerate(list_parts):  # 輪詢切好的一塊塊段落群
                str_part = str_part.strip()  # 字串變數，自己先削去前後隱形空白
                if str_part:  # 若內容非空無
                    list_lines = str_part.split('\n')  # 陣列變數，再將大段落沿換行切碎
                    str_title = list_lines[0].strip() if list_lines else f"Section {int_i}"  # 字串變數，以每一段的第一行做標題，不然就生 Section 系列編碼
                    str_body = '\n'.join(list_lines[1:]).strip() if len(list_lines) > 1 else ''  # 字串變數，將剩餘行數糊死作為本段內文
                    list_sections.append({  # 打包為標準物件追加
                        "title": str_title,  # 就定位
                        "content": str_body  # 回歸本位
                    })
        else:  # 若該文件純潔無瑕未見任何井號標記
            list_sections = [{  # 陣列變數，直接塞一包唯一字典硬吃
                "title": os.path.splitext(os.path.basename(str_file_path))[0],  # 從路徑去剝出原本的主檔名權當唯一標題
                "content": str_content  # 內文就是原本全文字
            }]
        
        dict_result = {  # 字典變數，向系統宣稱結果的標準格式
            "source_type": "txt",  # 告知後端為 txt
            "sections": list_sections if list_sections else [{"title": "Text Content", "content": "Empty file"}]  # 載入區塊或空檔字語
        }
        return dict_result  # 送出這個處理後的字典
    except Exception as err_txt:  # 處理那些檔案連開起來都失敗之絕望時刻
        dict_err = {  # 字典變數，組合包含 Exception 原文之報錯 JSON
            "source_type": "txt",  # 品種
            "sections": [{  # 建置錯誤區
                "title": "Error Reading Text File",  # 明講死因
                "content": f"Error: {str(err_txt)}"  # 附帶死因證據文字
            }]
        }
        return dict_err  # 打退堂鼓往上送